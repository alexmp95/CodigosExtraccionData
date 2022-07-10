from pptx import Presentation
X = Presentation("C:/Users/Alexis/Desktop/Pruebas/V2/data/Indicaciones Caso1.pptx") # Presentation object created
# Then file is opened in write mode
ftw_data = open("C:/Users/Alexis/Desktop/Pruebas/V2/data/Indicaciones Caso1.txt", "w", encoding='utf-8')
# write text from powerpoint
# file into .txt file
for slide in X.slides:

    for shape in slide.shapes:

        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:

                ftw_data.write(run.text)
ftw_data.close() # The file is closed